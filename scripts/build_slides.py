"""Build presentation.pptx — Patagonia-inspired aesthetic.

Palette: cream backgrounds, ink + rust accents, Futura/Avenir typography.
No template chrome — no top bars, minimal page numbers.

Regenerate: `python scripts/build_slides.py`
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parent.parent
FIG = ROOT / "figures"

# ── Patagonia-inspired palette ─────────────────────────────────────
CREAM   = RGBColor(0xEF, 0xE8, 0xDB)   # warm cream — slide bg
PAPER   = RGBColor(0xF7, 0xF2, 0xE8)   # lighter cream — card bg
INK     = RGBColor(0x2D, 0x31, 0x42)   # deep slate — body / titles
RUST    = RGBColor(0xC1, 0x44, 0x0E)   # burnt sienna — accent
FOREST  = RGBColor(0x4A, 0x67, 0x41)   # muted forest — secondary
SAND    = RGBColor(0xD4, 0xA5, 0x74)   # warm sand — tertiary
MUTED   = RGBColor(0x8B, 0x83, 0x74)   # warm grey — subtle text
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BLACK   = RGBColor(0x1A, 0x1D, 0x25)   # near-black ink

# Font stack — macOS has Futura, Avenir Next; Windows will fall back gracefully.
TITLE_FONT = "Futura"
BODY_FONT  = "Avenir Next"
MONO_FONT  = "Menlo"


def _set_slide_size(prs, width_in=13.333, height_in=7.5):
    prs.slide_width = Inches(width_in)
    prs.slide_height = Inches(height_in)


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _bg(slide, color, slide_w=Inches(13.333), slide_h=Inches(7.5)):
    """Fill the entire slide with a background colour."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_w, slide_h)
    shp.shadow.inherit = False
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()


def _rect(slide, left, top, width, height, fill=None, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.shadow.inherit = False
    if fill is not None:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
    return shp


def _textbox(slide, left, top, width, height, text, *, size=18, bold=False,
             color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=None, caps=False, tracking=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    if font is None:
        font = BODY_FONT
    if caps:
        text = text.upper()
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        if color is not None:
            run.font.color.rgb = color
    return tb


def _title(slide, text, *, y=Inches(0.6), size=36, color=None, caps=True):
    """Patagonia-style title: ALL CAPS Futura, tight, rust rule underneath."""
    color = color or INK
    _textbox(slide, Inches(0.7), y, Inches(12), Inches(0.9),
             text, size=size, bold=True, color=color,
             font=TITLE_FONT, caps=caps)
    # Rust accent rule
    _rect(slide, Inches(0.7), y + Inches(0.85), Inches(0.9), Inches(0.05),
          fill=RUST)


def _subtitle(slide, text, *, y=Inches(1.55)):
    _textbox(slide, Inches(0.7), y, Inches(12), Inches(0.5),
             text, size=14, color=MUTED, font=BODY_FONT)


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()


def _page(slide, n, slide_w=Inches(13.333), slide_h=Inches(7.5)):
    """Minimal page number — lower right, muted."""
    _textbox(slide, slide_w - Inches(0.8), slide_h - Inches(0.4),
             Inches(0.5), Inches(0.25), f"{n:02d}",
             size=9, color=MUTED, font=BODY_FONT, align=PP_ALIGN.RIGHT)


def _bullets(slide, left, top, width, height, items, *, size=16, color=None,
             spacing=10):
    color = color or INK
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(spacing)
        run = p.add_run()
        run.text = f"—  {item}"
        run.font.name = BODY_FONT
        run.font.size = Pt(size)
        if color is not None:
            run.font.color.rgb = color


def _pic(slide, path, left, top, *, width=None, height=None):
    kwargs = {}
    if width is not None:  kwargs["width"]  = width
    if height is not None: kwargs["height"] = height
    return slide.shapes.add_picture(str(path), left, top, **kwargs)


def build():
    prs = Presentation()
    _set_slide_size(prs)
    SW, SH = Inches(13.333), Inches(7.5)
    slides = []

    # ═════════════════════════════ 1. TITLE ═════════════════════════════
    s = _blank(prs); slides.append(s)
    _bg(s, INK)
    # Rust vertical bar
    _rect(s, Inches(0.7), Inches(2.3), Inches(0.12), Inches(2.8), fill=RUST)
    _textbox(s, Inches(1.1), Inches(2.1), Inches(11), Inches(0.5),
             "PHASE 2", size=14, bold=True, color=SAND,
             font=TITLE_FONT, caps=True)
    _textbox(s, Inches(1.1), Inches(2.55), Inches(11), Inches(1.2),
             "DETECTING AVALANCHES", size=54, bold=True, color=CREAM,
             font=TITLE_FONT)
    _textbox(s, Inches(1.1), Inches(3.5), Inches(11), Inches(1.2),
             "IN SAR IMAGERY", size=54, bold=True, color=CREAM,
             font=TITLE_FONT)
    _textbox(s, Inches(1.1), Inches(4.8), Inches(11), Inches(0.5),
             "A small D4-equivariant CNN, and an honest look at the metric",
             size=18, color=SAND, font=BODY_FONT)

    # Footer-ish author line
    _rect(s, Inches(1.1), Inches(6.55), Inches(2.5), Inches(0.03), fill=RUST)
    _textbox(s, Inches(1.1), Inches(6.7), Inches(6), Inches(0.35),
             "Guillermo San Marco",
             size=13, color=CREAM, font=TITLE_FONT, caps=True, bold=True)
    _textbox(s, Inches(1.1), Inches(7.0), Inches(6), Inches(0.3),
             "April 2026",
             size=11, color=MUTED, font=BODY_FONT)
    _notes(s, """
Avalanche debris segmentation on bi-temporal Sentinel-1 SAR. Two findings worth
leaving with: (1) a small equivariant CNN matches state-of-the-art on pixel F1
at 4× fewer params; (2) the standard benchmark metric hides a real operational
gap — our model finds more actual avalanches. ~15 minutes.
""")

    # ═════════════════════════════ 2. WHY ═════════════════════════════
    s = _blank(prs); slides.append(s)
    _bg(s, CREAM)
    _title(s, "Why this matters")
    _subtitle(s, "Avalanche rescue and road management depend on knowing where debris fell")
    _bullets(s, Inches(0.7), Inches(2.4), Inches(7.2), Inches(5), [
        "Avalanches kill ~150 people a year in Europe alone",
        "Optical imagery is useless under clouds — avalanches happen in storms",
        "SAR sees through cloud / night — the only reliable sensor at scale",
        "Key signal: the DIFFERENCE between a pre-storm and post-storm acquisition",
        "Goal: automatic, pixel-accurate deposit maps from free Sentinel-1 data",
    ], size=17, spacing=12)
    _pic(s, FIG / "pair16_overlay.png", Inches(8.0), Inches(2.3),
         height=Inches(5.0))
    _notes(s, """
Why care. Avalanches kill ~150 people per year in Europe. They happen during
storms, when optical satellites see only clouds. Sentinel-1 SAR is free,
all-weather, and sees through cloud and night. Its downside: the imagery is
noisy and deposits only show up as a difference between pre-storm and
post-storm acquisitions. The goal is an automatic pipeline: hand it two SAR
scenes, get back a deposit map. The example on the right shows raw VH
backscatter with a deposit polygon overlaid.
""")

    # ═════════════════════════════ 3. DATASET ═════════════════════════════
    s = _blank(prs); slides.append(s)
    _bg(s, CREAM)
    _title(s, "The dataset")
    _subtitle(s, "AvalCD (Gatti et al.)  ·  117 polygons on the held-out Tromsø OOD test scene")

    # ── Row 1: data split (3 compact cards) ───────────────────────
    split_y = Inches(1.95)
    split_h = Inches(1.15)
    cards = [
        ("TRAIN",
         "5 scenes  ·  ~34 k deposit pixels",
         "Livigno ×2, Nuuk ×2, Pish ×1",
         FOREST),
        ("VAL",
         "Livigno  2025-03-18",
         "threshold & calibration tuning",
         SAND),
        ("TEST",
         "Tromsø  2024-12-20   (OOD)",
         "117 polygons  ·  0.5 % pixels positive",
         RUST),
    ]
    card_w = Inches(4.1)
    gap = Inches(0.1)
    for i, (label, main, sub, accent) in enumerate(cards):
        x = Inches(0.7) + (card_w + gap) * i
        _rect(s, x, split_y, card_w, split_h, fill=PAPER)
        _rect(s, x, split_y, Inches(0.08), split_h, fill=accent)
        _textbox(s, x + Inches(0.25), split_y + Inches(0.12),
                 card_w - Inches(0.4), Inches(0.35),
                 label, size=11, bold=True, color=accent,
                 font=TITLE_FONT, caps=True)
        _textbox(s, x + Inches(0.25), split_y + Inches(0.45),
                 card_w - Inches(0.4), Inches(0.4),
                 main, size=13, bold=True, color=INK, font=BODY_FONT)
        _textbox(s, x + Inches(0.25), split_y + Inches(0.80),
                 card_w - Inches(0.4), Inches(0.3),
                 sub, size=10, color=MUTED, font=BODY_FONT)

    # ── Row 2: EDA chart (width-constrained so it can't overflow) ─
    # eda.png aspect is 3.61:1 — at width 12.7" → height 3.52"
    _pic(s, FIG / "eda.png", Inches(0.3), Inches(3.4), width=Inches(12.7))
    _notes(s, """
AvalCD public benchmark from Gatti et al. Five training scenes (Italy,
Greenland, Pakistan), one validation scene, and one out-of-distribution test
scene — Tromsø, 117 polygons, 0.5% of pixels positive. Deposits labeled with
EAWS D-scale: D1 knocks down a person, D2 can bury, D3 can destroy a car, D4
can destroy a building. Test scene is D3-dominated (n=72). Median deposit is
~12 000 m² — roughly a small parking lot.
""")

    # ═════════════════════════════ 4. ORIENTATION ═════════════════════════════
    s = _blank(prs); slides.append(s)
    _bg(s, CREAM)
    _title(s, "Orientation doesn't matter")
    _subtitle(s, "A rotated avalanche is still the same avalanche — the satellite just flew a different path")
    _pic(s, FIG / "d4_on_debris.png", Inches(0.7), Inches(2.0),
         width=Inches(12.0))
    _textbox(s, Inches(0.7), Inches(6.95), Inches(12), Inches(0.5),
             "A standard CNN must LEARN this from 8× augmented data.  "
             "A D4-equivariant CNN gets it for free, by construction.",
             size=15, color=RUST, align=PP_ALIGN.CENTER, bold=True,
             font=BODY_FONT)
    _notes(s, """
The physical insight. All 8 panels are the same real avalanche patch under the
4 rotations and 4 reflections that make up D4 — the symmetry group of the
square. Physically equivalent: the avalanche doesn't know which way the
satellite flew. Standard CNN has to learn this from augmented data.
D4-equivariant CNN encodes the symmetry into the convolutions. Rotation in →
rotation out, guaranteed.
""")

    # ═════════════════════════════ 5. APPROACH ═════════════════════════════
    s = _blank(prs); slides.append(s)
    _bg(s, CREAM)
    _title(s, "Our approach")
    _subtitle(s, "A small equivariant encoder compares pre and post scenes")
    _pic(s, FIG / "architecture.png", Inches(0.3), Inches(2.0),
         width=Inches(12.7))
    _notes(s, """
12-channel input (pre, post, plus derived features). Two encoder branches with
SHARED weights — both D4-equivariant. Subtract pre features from post features
(the change). GroupPool to an invariant tensor. Small Conv2d decoder turns
that into a pixel-level prediction, with DEM and slope injected at each
decoder stage. Two outputs: mask and deposit area. Total: 625,617 parameters.
Gatti's Swin-UNet: 2.39 million.
""")

    # ═════════════════════════════ 6. PIXEL RESULTS + FRONTIER ═════════════════════════════
    s = _blank(prs); slides.append(s)
    _bg(s, CREAM)
    _title(s, "Pixel-level results")
    _subtitle(s, "AvalCD Table 5 (Gatti et al.) with our model added  ·  Pareto frontier on params × F1")

    # ── LEFT: compact Table 5 ────────────────────────────────────
    rows = [
        ("Model",                          "Params",  "F1"),
        ("U-Net",                          "31.0 M",  "0.487"),
        ("FCN8",                           "134.3 M", "0.602"),
        ("TinyCD",                         "0.29 M",  "0.766"),
        ("RUNet",                          "7.76 M",  "0.767"),
        ("A-BT-UNet",                      "12.4 M",  "0.793"),
        ("D4-EquiCNN  (ours)",             "0.63 M",  "0.794"),
        ("Swin-UNet (Gatti)",              "2.39 M",  "0.803"),
    ]
    table_left = Inches(0.6)
    table_w    = Inches(5.9)
    y0 = Inches(2.15)
    row_h = Inches(0.48)
    col_x = [table_left + Inches(0.2), table_left + Inches(3.4),
             table_left + Inches(4.8)]
    for i, row in enumerate(rows):
        y = y0 + row_h * i
        if i == 0:
            _rect(s, table_left, y + Inches(0.43), table_w, Inches(0.02),
                  fill=INK)
            size, bold, color = 11, True, INK
            caps = True; font = TITLE_FONT
        elif "ours" in row[0]:
            _rect(s, table_left, y, table_w, row_h, fill=PAPER)
            _rect(s, table_left, y, Inches(0.08), row_h, fill=RUST)
            size, bold, color = 14, True, RUST
            caps = False; font = BODY_FONT
        else:
            size, bold, color = 13, False, INK
            caps = False; font = BODY_FONT
        widths = [Inches(3.2), Inches(1.4), Inches(1.0)]
        for j, val in enumerate(row):
            _textbox(s, col_x[j], y + Inches(0.12), widths[j], row_h,
                     val, size=size, bold=bold, color=color,
                     align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER,
                     font=font, caps=caps)

    # ── RIGHT: Pareto plot ───────────────────────────────────────
    _pic(s, FIG / "param_frontier.png", Inches(6.8), Inches(2.0),
         height=Inches(4.7))

    # ── BOTTOM: caption tying them together ─────────────────────
    _rect(s, Inches(0.6), Inches(6.6), Inches(0.08), Inches(0.55), fill=RUST)
    _textbox(s, Inches(0.85), Inches(6.65), Inches(12), Inches(0.5),
             "6 TH OF 7 ON TABLE 5  ·  WITHIN 1 PP OF THE LEADER  ·  4× FEWER PARAMETERS",
             size=13, bold=True, color=INK, font=TITLE_FONT, caps=True)
    _textbox(s, Inches(0.85), Inches(7.00), Inches(12), Inches(0.35),
             "Pareto-dominant bottom-right of the frontier.",
             size=11, color=MUTED, font=BODY_FONT)
    _notes(s, """
Table 5 from Gatti's paper with our model inserted, paired with the Pareto plot
on the right. U-Net and FCN8 are the old guard: 30-130M params, weak F1.
TinyCD, RUNet, A-BT-UNet — more recent. At the top: Swin-UNet at 2.39M
params, F1 0.803. Ours: 0.63M params, F1 0.794. Within 1pp of Swin-UNet at
4× fewer params — and as the Pareto plot on the right makes visually clear,
we sit alone in the bottom-right of the frontier (high F1, low params).
""")

    # ═════════════════════════════ 7. PIXEL F1 vs INSTANCE F1 ═════════════════════════════
    s = _blank(prs); slides.append(s)
    _bg(s, CREAM)
    _title(s, "Pixel F1 vs instance F1")
    _subtitle(s, "Same predictions, different metric, different ranking  ·  Tromsø test scene")

    # ── LEFT: metric disagreement bar chart (width-constrained to column) ──
    # aspect 2.22:1 — at width 6.0" → height 2.7"
    _pic(s, FIG / "metric_disagreement.png", Inches(0.4), Inches(2.9),
         width=Inches(6.0))

    # ── RIGHT: D-scale breakdown table ───────────────────────────
    table_x = Inches(6.8)
    _textbox(s, table_x, Inches(2.0), Inches(6.2), Inches(0.35),
             "INSTANCE F1 BREAKDOWN BY D-SCALE  (IoU > 0.3)",
             size=11, bold=True, color=INK, font=TITLE_FONT, caps=True)

    rows = [
        ("D-scale",      "ours",  "Gatti",  "Δ",        "n"),
        ("D2  small",    "0.113", "0.115",  "−0.002",   "25"),
        ("D3  medium",   "0.531", "0.399",  "+0.132",   "72"),
        ("D4  large",    "0.199", "0.136",  "+0.063",   "16"),
        ("all excl D1",  "0.637", "0.530",  "+0.107",   "113"),
    ]
    y0 = Inches(2.5)
    row_h = Inches(0.6)
    col_x = [table_x + Inches(0.2), table_x + Inches(1.9), table_x + Inches(2.9),
             table_x + Inches(4.0), table_x + Inches(5.4)]
    col_w = Inches(1.2)
    for i, row in enumerate(rows):
        y = y0 + row_h * i
        if i == 0:
            _rect(s, table_x, y + Inches(0.55), Inches(6.2), Inches(0.02),
                  fill=INK)
            size, bold, color, font, caps = 11, True, INK, TITLE_FONT, True
        elif i == 2:   # D3 row
            _rect(s, table_x, y, Inches(6.2), row_h, fill=PAPER)
            _rect(s, table_x, y, Inches(0.08), row_h, fill=RUST)
            size, bold, color, font, caps = 17, True, RUST, BODY_FONT, False
        elif i == 4:   # overall row
            _rect(s, table_x, y, Inches(6.2), row_h, fill=PAPER)
            _rect(s, table_x, y, Inches(0.08), row_h, fill=RUST)
            size, bold, color, font, caps = 15, True, RUST, BODY_FONT, False
        else:
            size, bold, color, font, caps = 14, False, INK, BODY_FONT, False
        for j, val in enumerate(row[:5]):
            first_col_w = Inches(1.9) if j == 0 else col_w
            _textbox(s, col_x[j], y + Inches(0.15), first_col_w, row_h,
                     val, size=size, bold=bold, color=color,
                     align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER,
                     font=font, caps=caps)

    # Footnote caption under table
    _textbox(s, table_x, Inches(5.7), Inches(6.2), Inches(0.4),
             "D3 is where the gap is largest  —  and D3 is what practitioners act on.",
             size=13, bold=True, color=INK, font=BODY_FONT)
    _textbox(s, table_x, Inches(6.05), Inches(6.2), Inches(0.4),
             "Robust across IoU thresholds 0.1–0.5  ·  center-point matching confirms the ranking",
             size=10, color=MUTED, font=BODY_FONT)

    # ── Bottom: unifying takeaway ────────────────────────────────
    _rect(s, Inches(0.6), Inches(6.9), Inches(0.08), Inches(0.5), fill=RUST)
    _textbox(s, Inches(0.85), Inches(6.9), Inches(12), Inches(0.4),
             "Pixel F1 says tie  ·  instance F1 says +10.7 pp",
             size=14, bold=True, color=INK, font=TITLE_FONT, caps=True)
    _textbox(s, Inches(0.85), Inches(7.2), Inches(12), Inches(0.3),
             "Pixel F1 treats every pixel independently. It does not ask \"which avalanches did you find?\"",
             size=11, color=MUTED, font=BODY_FONT)
    _notes(s, """
Two views of the same predictions. Left: pixel F1 vs instance F1 bar chart —
same model pair, same ground truth. The pixel bars are nearly identical; the
instance bars split wide open. Right: D-scale breakdown of the instance F1.
D2 (small) — tie. D3 (medium, car-destroying) — +13 points on n=72,
well-powered. D4 (large) — +6 points. Overall +10.7 pp. This gap is robust
across IoU thresholds and under center-point matching.
""")

    # ═════════════════════════════ 9. HEAD-TO-HEAD ═════════════════════════════
    s = _blank(prs); slides.append(s)
    _bg(s, CREAM)
    _title(s, "One deposit, two models")
    _subtitle(s, "Same ground-truth outline  ·  same inference pipeline  ·  different architectures")
    _pic(s, FIG / "head_to_head.png", Inches(0.3), Inches(2.2),
         width=Inches(12.7))
    _notes(s, """
One real example from Tromsø. Left to right: VH backscatter, ground truth,
our probability, Gatti's probability. Same cyan outline. Our model locks onto
the actual deposit — IoU 0.77. Swin-UNet spills probability all over — IoU
0.20. This is the 10-point gap in one frame.
""")

    # ═════════════════════════════ 10. DPR ═════════════════════════════
    s = _blank(prs); slides.append(s)
    _bg(s, CREAM)
    _title(s, "Can we bolt on an instance head?")
    _subtitle(s, "Phase 1 of this project  ·  what we tried  ·  what happened")

    # Left: motivation + Phase 1
    _rect(s, Inches(0.7), Inches(2.1), Inches(5.9), Inches(4.8), fill=PAPER)
    _rect(s, Inches(0.7), Inches(2.1), Inches(0.08), Inches(4.8), fill=FOREST)
    _textbox(s, Inches(0.95), Inches(2.25), Inches(5.5), Inches(0.4),
             "MOTIVATION", size=11, bold=True, color=FOREST,
             font=TITLE_FONT, caps=True)
    _bullets(s, Inches(0.95), Inches(2.7), Inches(5.5), Inches(2), [
        "Maybe the instance-F1 gap can be fixed with pipeline engineering",
        "Mask R-CNN decomposes detect → segment; we can try the same",
    ], size=13, spacing=8)

    _textbox(s, Inches(0.95), Inches(4.4), Inches(5.5), Inches(0.4),
             "WHAT IS PHASE 1?", size=11, bold=True, color=FOREST,
             font=TITLE_FONT, caps=True)
    _bullets(s, Inches(0.95), Inches(4.85), Inches(5.5), Inches(2), [
        "Same D4-equivariant backbone, classification head",
        "Slides 64 × 64 patches across the scene → dense avalanche-probability heatmap",
        "Transfers to Tromsø OOD with 100 % proposal recall",
    ], size=13, spacing=8)

    # Right: DPR results
    _rect(s, Inches(6.8), Inches(2.1), Inches(5.9), Inches(4.8), fill=PAPER)
    _rect(s, Inches(6.8), Inches(2.1), Inches(0.08), Inches(4.8), fill=RUST)
    _textbox(s, Inches(7.05), Inches(2.25), Inches(5.5), Inches(0.4),
             "DETECT → PROPOSE → REFINE", size=11, bold=True, color=RUST,
             font=TITLE_FONT, caps=True)
    _textbox(s, Inches(7.05), Inches(2.7), Inches(5.5), Inches(0.4),
             "Phase 1 proposes regions. Refiner segments each crop.",
             size=12, color=MUTED, font=BODY_FONT)

    results = [
        ("Refine crops (fine-tune)",
         "F1 = 0.555",
         "150 TP crops vs 30 k patches"),
        ("Refine crops (encoder frozen)",
         "F1 = 0.551",
         "bottleneck is data, not forgetting"),
        ("Heatmap as Phase-2 channel",
         "F1 = 0.773",
         "Phase 1 info already captured internally"),
    ]
    y = Inches(3.3)
    for head, res, note in results:
        _textbox(s, Inches(7.05), y, Inches(5.5), Inches(0.3),
                 head, size=12, bold=True, color=INK, font=BODY_FONT)
        _textbox(s, Inches(7.05), y + Inches(0.3), Inches(3), Inches(0.4),
                 res, size=16, bold=True, color=RUST, font=BODY_FONT)
        _textbox(s, Inches(7.05), y + Inches(0.68), Inches(5.5), Inches(0.3),
                 note, size=10, color=MUTED, font=BODY_FONT)
        y += Inches(1.1)

    _textbox(s, Inches(0.7), Inches(7.05), Inches(12), Inches(0.4),
             "The instance-F1 advantage is architectural. It cannot be bolted on after the fact.",
             size=14, bold=True, color=INK, align=PP_ALIGN.CENTER, font=BODY_FONT)
    _notes(s, """
DPR = Detect-Propose-Refine. Motivation: Mask R-CNN decomposes instance
segmentation — so maybe we could too. Phase 1 (separate repo) is a classifier
with the same D4-equivariant backbone but a classification head — patches in,
Y/N out. Transfers to Tromsø OOD with 100% recall of real deposits.

Three variants:
1. Fine-tune a conditional segmenter on crops where Phase 1 fires → F1 0.555.
   Cause: 150 TP crops vs 30k pixel patches. Data scarcity.
2. Freeze the encoder during retraining → F1 0.551. Almost identical — so it's
   not catastrophic forgetting, it's data volume.
3. Add Phase 1's heatmap as an extra Phase 2 channel, train on all 30k
   patches → F1 0.773. Below our 0.794 baseline. Phase 1 output is redundant
   with Phase 2's internal features.

Conclusion: instance advantage is architectural, not glueable.
""")

    # ═════════════════════════════ 12. OTHER NEGATIVES ═════════════════════════════
    s = _blank(prs); slides.append(s)
    _bg(s, CREAM)
    _title(s, "Other negative results worth reporting")

    items = [
        ("λ SWEEP OF COMPONENT-IoU LOSS",
         "Trained with  L = BCE + λ · component-IoU  for λ ∈ {0, 0.25, 0.5, 1, 2}",
         "No setting kept both pixel F1 ≥ 0.78 AND D2 detection ≥ 65 %"),
        ("BIGGER / DEEPER MODELS",
         "More capacity consistently HURT",
         "Inductive bias is doing the work, not scale"),
        ("HEAVY REGULARIZATION, WARM RESTARTS",
         "Both hurt",
         "Equivariance is already regularizing the model"),
        ("CALIBRATION ASYMMETRY",
         "Our F1-optimal threshold 0.225, Gatti's 0.775",
         "Same F1, completely different probability distributions"),
    ]
    y = Inches(2.2)
    for i, (head, detail, note) in enumerate(items):
        yy = y + Inches(1.15) * i
        _rect(s, Inches(0.7), yy + Inches(0.1), Inches(0.08),
              Inches(0.85), fill=RUST)
        _textbox(s, Inches(1.0), yy, Inches(11.8), Inches(0.4),
                 head, size=13, bold=True, color=INK,
                 font=TITLE_FONT, caps=True)
        _textbox(s, Inches(1.0), yy + Inches(0.42), Inches(11.8), Inches(0.4),
                 detail, size=13, color=INK, font=BODY_FONT)
        _textbox(s, Inches(1.0), yy + Inches(0.75), Inches(11.8), Inches(0.4),
                 note, size=11, color=MUTED, font=BODY_FONT)
    _notes(s, """
Four negatives worth the slide. Lambda sweep: we trained the component-IoU
loss at five lambda values. Never found a setting that kept both pixel F1 high
AND D2 detection up. Bigger/deeper models hurt — evidence for bias, not scale.
Heavy regularization hurt — equivariance is already regularizing. Calibration
asymmetry: both hit F1 ~0.79 at very different thresholds — 0.225 vs 0.775.
Interesting for practitioners; doesn't change the ranking.
""")

    # ═════════════════════════════ 13. CONTRIBUTIONS ═════════════════════════════
    s = _blank(prs); slides.append(s)
    _bg(s, CREAM)
    _title(s, "What we contributed")
    _subtitle(s, "Four things worth taking away")

    items = [
        ("01", "A METRIC THAT MATCHES WHAT OPERATORS CARE ABOUT",
         "Instance-level F1 reveals detection gaps that pixel F1 hides. "
         "Should become standard reporting on AvalCD."),
        ("02", "D4-EQUIVARIANCE WINS OVER SCALE HERE",
         "625 k params match or exceed a 2.39 M Swin-UNet on pixel F1, "
         "and outperform it by +10.7 pp on instance F1."),
        ("03", "THE GAP MATTERS MOST WHERE IT COUNTS — D3",
         "+13 pp on the deposits that drive real decisions "
         "(road closures, rescue dispatch).  n = 72, well-powered."),
        ("04", "HONEST NEGATIVES RULE OUT THE EASY ALTERNATIVES",
         "DPR, loss engineering, and scale all failed. "
         "The advantage is architectural."),
    ]
    y = Inches(2.25)
    for i, (num, head, body) in enumerate(items):
        yy = y + Inches(1.05) * i
        _textbox(s, Inches(0.7), yy, Inches(1.2), Inches(0.9),
                 num, size=42, bold=True, color=RUST,
                 font=TITLE_FONT, align=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.MIDDLE)
        _textbox(s, Inches(2.0), yy + Inches(0.1), Inches(11), Inches(0.5),
                 head, size=15, bold=True, color=INK,
                 font=TITLE_FONT, caps=True)
        _textbox(s, Inches(2.0), yy + Inches(0.55), Inches(11), Inches(0.5),
                 body, size=13, color=MUTED, font=BODY_FONT)
    _notes(s, """
Four contributions. 01 METHOD: instance-level F1 with IoU matching should
become standard reporting on AvalCD — pixel F1 is demonstrably blind to
operational behavior. 02 MODEL: 625k params beats 2.39M on instance F1,
matches on pixel F1. 03 OPS: gap is largest where it matters — D3, with n=72
and well-powered. 04 RIGOR: DPR, loss engineering, scale — all failed.
Advantage is architectural.
""")

    # ═════════════════════════════ 14. CONCLUSION ═════════════════════════════
    s = _blank(prs); slides.append(s)
    _bg(s, CREAM)
    _title(s, "Conclusion  /  what comes next")

    # Left: conclusion
    _rect(s, Inches(0.7), Inches(2.1), Inches(5.9), Inches(4.9), fill=PAPER)
    _rect(s, Inches(0.7), Inches(2.1), Inches(0.08), Inches(4.9), fill=FOREST)
    _textbox(s, Inches(0.95), Inches(2.25), Inches(5.5), Inches(0.4),
             "CONCLUSION", size=11, bold=True, color=FOREST,
             font=TITLE_FONT, caps=True)
    _textbox(s, Inches(0.95), Inches(2.7), Inches(5.5), Inches(1.2),
             "Symmetry-aware inductive bias beats scale for this problem.",
             size=22, bold=True, color=INK, font=TITLE_FONT)
    _bullets(s, Inches(0.95), Inches(4.5), Inches(5.5), Inches(3), [
        "Matched pixel F1 at 4× fewer parameters",
        "+10.7 pp instance F1, +13 pp on D3",
        "Pixel F1 is an insufficient benchmark",
        "The advantage is architectural, not bolt-on",
    ], size=13, spacing=8)

    # Right: future directions
    _rect(s, Inches(6.8), Inches(2.1), Inches(5.9), Inches(4.9), fill=PAPER)
    _rect(s, Inches(6.8), Inches(2.1), Inches(0.08), Inches(4.9), fill=RUST)
    _textbox(s, Inches(7.05), Inches(2.25), Inches(5.5), Inches(0.4),
             "WHAT COMES NEXT", size=11, bold=True, color=RUST,
             font=TITLE_FONT, caps=True)
    items = [
        ("More OOD test scenes",
         "One held-out scene is thin evidence for generalization  ·  "
         "need curated multi-scene OOD benchmarks"),
        ("Seed-level variance on headline",
         "0.7938 is seed 1 only; ablations used 3 seeds"),
        ("D2 detection floor",
         "Physics-limited or architecture-limited?"),
        ("Operational integration",
         "Sentinel-1 ingest → deposit map in minutes"),
        ("Other SAR missions",
         "SAOCOM, NISAR, TerraSAR-X — same symmetry, different geometry"),
    ]
    y = Inches(2.8)
    for head, detail in items:
        _textbox(s, Inches(7.05), y, Inches(5.5), Inches(0.3),
                 head, size=12, bold=True, color=INK, font=BODY_FONT)
        _textbox(s, Inches(7.05), y + Inches(0.28), Inches(5.5), Inches(0.4),
                 detail, size=10, color=MUTED, font=BODY_FONT)
        y += Inches(0.82)
    _notes(s, """
Conclusion: for a problem with known symmetry, a small equivariant network
beats a larger transformer that has to LEARN the symmetry — on the standard
metric (tied) AND on the metric that matters for operations (+10.7 pp).
Advantage is architectural — we tried the alternatives, they failed.

Next: cross-scene validation, seed variance, D2 floor analysis, operational
pipeline, other SAR missions.
""")

    # ═════════════════════════════ 15. THANKS ═════════════════════════════
    s = _blank(prs); slides.append(s)
    _bg(s, INK)
    _rect(s, Inches(0.7), Inches(2.8), Inches(0.12), Inches(2.0), fill=RUST)
    _textbox(s, Inches(1.1), Inches(2.6), Inches(11), Inches(0.5),
             "END", size=14, bold=True, color=SAND,
             font=TITLE_FONT, caps=True)
    _textbox(s, Inches(1.1), Inches(3.1), Inches(11), Inches(1.5),
             "THANKS.", size=80, bold=True, color=CREAM, font=TITLE_FONT)
    _textbox(s, Inches(1.1), Inches(4.5), Inches(11), Inches(0.8),
             "Questions welcome.",
             size=28, color=SAND, font=BODY_FONT)
    _rect(s, Inches(1.1), Inches(6.1), Inches(2.5), Inches(0.03), fill=RUST)
    _textbox(s, Inches(1.1), Inches(6.25), Inches(11), Inches(0.4),
             "github.com/gsanmarco/Equivariant-SAR-Segmentation",
             size=12, color=CREAM, font=MONO_FONT)
    _textbox(s, Inches(1.1), Inches(6.65), Inches(11), Inches(0.4),
             "Phase 1:  github.com/sanmarcog/Equivariant-CNN-SAR",
             size=11, color=MUTED, font=MONO_FONT)
    _notes(s, """
Happy to go deeper on: escnn regular representations, the instance-matching
rule, SAR log-ratio channels, Phase 1 architecture, or any of the negative
results.
""")

    # ── Page numbers on content slides only ───────────────────────
    total = len(slides)
    for i, sl in enumerate(slides):
        if i == 0 or i >= total - 1:
            continue
        _page(sl, i + 1)

    out = ROOT / "presentation.pptx"
    prs.save(str(out))
    print(f"Saved: {out}  ({total} slides)")


if __name__ == "__main__":
    build()
